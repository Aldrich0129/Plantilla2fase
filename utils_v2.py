"""
Librería común para el sistema de automatización de plantillas
Funciones compartidas entre Fase 1 y Fase 2
VERSIÓN 2: Con soporte para contexto de variables y patrones combinados
"""

import re
from datetime import datetime
from typing import Any, Dict, List, Set, Tuple
from docx import Document
from docx.shared import RGBColor
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt
import yaml


class PatternDetector:
    """Detecta patrones de variables en documentos"""

    # Patrones de texto predefinidos
    TEXT_PATTERNS = {
        'llaves_simples': r'\{([^}]+)\}',
        'llaves_dobles': r'\{\{([^}]+)\}\}',
        'corchetes_simples': r'\[([^\]]+)\]',
        'corchetes_dobles': r'\[\[([^\]]+)\]\]',
    }

    # Patrones de fecha con "de" (para detectar como variable única)
    # ORDEN DE PRIORIDAD: más largo a más corto
    DATE_WITH_DE_PATTERNS = [
        r'(d[ií]a\s+de\s+mes\s+de\s+a[ñn]o)',  # día de mes de año (PRIORIDAD 1)
        r'(d[ií]a\s+de\s+mes)',  # día de mes (PRIORIDAD 2)
        r'(mes\s+de\s+a[ñn]o)',  # mes de año (PRIORIDAD 3)
    ]
    
    @staticmethod
    def detect_colors_in_docx(doc: Document) -> Dict[str, Set[str]]:
        """Detecta todos los colores usados en un documento Word"""
        colors = {
            'text_colors': set(),
            'highlight_colors': set()
        }
        
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                # Color de texto
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    colors['text_colors'].add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                
                # Color de subrayado/resaltado
                if run.font.highlight_color:
                    colors['highlight_colors'].add(str(run.font.highlight_color))
        
        # También revisar tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run.font.color and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                colors['text_colors'].add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
        
        return colors
    
    @staticmethod
    def detect_colors_in_pptx(prs: Presentation) -> Dict[str, Set[str]]:
        """
        Detecta todos los colores usados en una presentación PowerPoint.
        Maneja diferentes tipos de colores: RGB directo, colores de tema, etc.
        """
        colors = {
            'text_colors': set(),
            'highlight_colors': set()
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame is None:
                        continue

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                # Verificar si el run tiene color
                                if not run.font.color:
                                    continue

                                color_obj = run.font.color
                                rgb_value = None

                                # Método 1: Intentar obtener RGB directo
                                if hasattr(color_obj, 'rgb'):
                                    try:
                                        rgb_value = color_obj.rgb
                                    except (AttributeError, TypeError):
                                        pass

                                # Método 2: Si es color de tema, intentar obtener RGB del tema
                                if rgb_value is None and hasattr(color_obj, 'theme_color'):
                                    try:
                                        # Los colores de tema no se pueden convertir fácilmente
                                        # Se omiten por ahora
                                        pass
                                    except (AttributeError, TypeError):
                                        pass

                                # Si obtuvimos un valor RGB válido, agregarlo
                                if rgb_value:
                                    colors['text_colors'].add(f"#{rgb_value[0]:02x}{rgb_value[1]:02x}{rgb_value[2]:02x}")

                            except Exception:
                                # Si hay cualquier otro error, simplemente continuar
                                # Esto evita que el programa se detenga por colores no soportados
                                pass
                except Exception:
                    # Los shapes sin text frame pueden lanzar errores si se accede directamente.
                    # Los ignoramos para evitar que falle la detección en PPTX complejos.
                    continue
        
        return colors
    
    @staticmethod
    def is_substring_of_any(text: str, text_list: List[str]) -> bool:
        """
        Verifica si un texto está contenido dentro de cualquier texto de la lista.
        Comparación case-insensitive.
        """
        text_lower = text.lower().strip()
        for other in text_list:
            other_lower = other.lower().strip()
            if text_lower != other_lower and text_lower in other_lower:
                return True
        return False

    @staticmethod
    def detect_date_with_de_patterns(text: str) -> List[str]:
        """
        Detecta patrones de fecha con 'de' como una variable única.
        Por ejemplo: 'día de mes de año', 'día de mes', 'mes de año'
        Respeta la prioridad: detecta primero los más largos y excluye subconjuntos.
        """
        detected_dates = []
        text_lower = text.lower()

        # Detectar en orden de prioridad (más largo a más corto)
        for pattern in PatternDetector.DATE_WITH_DE_PATTERNS:
            matches = re.finditer(pattern, text_lower, re.IGNORECASE)
            for match in matches:
                found_text = match.group(1)
                # Solo agregar si no es substring de algo ya detectado
                if not PatternDetector.is_substring_of_any(found_text, detected_dates):
                    detected_dates.append(found_text)

        return detected_dates  # Mantener orden de prioridad

    @staticmethod
    def extract_variables_by_pattern(text: str, pattern: str) -> List[str]:
        """Extrae variables según un patrón regex"""
        if pattern in PatternDetector.TEXT_PATTERNS:
            regex = PatternDetector.TEXT_PATTERNS[pattern]
        else:
            regex = pattern

        matches = re.findall(regex, text)
        return list(set(matches))  # Eliminar duplicados
    
    @staticmethod
    def extract_variables_by_color(doc: Document, color: str, color_type: str = 'text') -> List[Tuple[str, Any]]:
        """
        Extrae texto con un color específico, agrupando runs consecutivos del mismo color.
        
        Mejora clave: Si varios runs consecutivos tienen el mismo formato, 
        se agrupan en UNA sola variable.
        """
        variables = []
        
        # Procesar párrafos
        for paragraph in doc.paragraphs:
            grouped_vars = PatternDetector._group_consecutive_colored_runs(
                paragraph.runs, color, color_type
            )
            variables.extend(grouped_vars)
        
        # Procesar tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        grouped_vars = PatternDetector._group_consecutive_colored_runs(
                            paragraph.runs, color, color_type
                        )
                        variables.extend(grouped_vars)
        
        return variables
    
    @staticmethod
    def _group_consecutive_colored_runs(runs: list, color: str, color_type: str) -> List[Tuple[str, Any]]:
        """
        Agrupa runs consecutivos con el mismo color en una sola variable.
        
        Args:
            runs: Lista de runs del párrafo
            color: Color a buscar (hex para text, string para highlight)
            color_type: 'text' o 'highlight'
            
        Returns:
            Lista de tuplas (texto_agrupado, primer_run)
        """
        variables = []
        current_text = ""
        first_run = None
        
        for run in runs:
            run_matches = False
            
            # Verificar si el run tiene el color buscado
            if color_type == 'text' and run.font.color and run.font.color.rgb:
                rgb = run.font.color.rgb
                run_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                run_matches = (run_color == color)
            elif color_type == 'highlight' and run.font.highlight_color:
                run_matches = (str(run.font.highlight_color) == color)
            
            if run_matches:
                # Este run tiene el color buscado
                if first_run is None:
                    first_run = run
                current_text += run.text
            else:
                # Este run NO tiene el color, guardar lo acumulado
                if current_text.strip():
                    variables.append((current_text.strip(), first_run))
                current_text = ""
                first_run = None
        
        # Guardar el último grupo si existe
        if current_text.strip():
            variables.append((current_text.strip(), first_run))
        
        return variables
    
    @staticmethod
    def detect_list_options(text: str) -> List[str]:
        """
        Detecta si un texto contiene opciones separadas por delimitadores.
        Soporta: / | ,
        """
        delimiters = ['/', '|', ',']
        
        for delim in delimiters:
            if delim in text:
                parts = [p.strip() for p in text.split(delim) if p.strip()]
                if len(parts) >= 2:  # Al menos 2 opciones
                    return parts
        
        return []
    
    @staticmethod
    def extract_variable_context(doc: Document, variable_text: str, context_chars: int = 20) -> List[Dict[str, str]]:
        """
        Extrae el contexto de una variable en el documento.
        
        Args:
            doc: Documento Word
            variable_text: Texto de la variable a buscar
            context_chars: Cantidad de caracteres de contexto antes/después
            
        Returns:
            Lista de diccionarios con 'before', 'variable', 'after', 'location'
        """
        contexts = []
        location_counter = 1
        
        # Buscar en párrafos
        for paragraph in doc.paragraphs:
            full_text = paragraph.text
            contexts.extend(PatternDetector._extract_context_from_text(
                full_text, variable_text, context_chars, f"Párrafo {location_counter}"
            ))
            if variable_text in full_text:
                location_counter += 1
        
        # Buscar en tablas
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    for paragraph in cell.paragraphs:
                        full_text = paragraph.text
                        contexts.extend(PatternDetector._extract_context_from_text(
                            full_text, variable_text, context_chars, 
                            f"Tabla {table_idx+1}, Fila {row_idx+1}, Celda {cell_idx+1}"
                        ))
        
        # Buscar en encabezados y pies
        for section_idx, section in enumerate(doc.sections):
            if section.header:
                for paragraph in section.header.paragraphs:
                    full_text = paragraph.text
                    contexts.extend(PatternDetector._extract_context_from_text(
                        full_text, variable_text, context_chars, f"Encabezado {section_idx+1}"
                    ))
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    full_text = paragraph.text
                    contexts.extend(PatternDetector._extract_context_from_text(
                        full_text, variable_text, context_chars, f"Pie de página {section_idx+1}"
                    ))
        
        return contexts
    
    @staticmethod
    def extract_variable_context_pptx(prs: Presentation, variable_text: str, context_chars: int = 20) -> List[Dict[str, str]]:
        """
        Extrae el contexto de una variable en una presentación PowerPoint.
        
        Args:
            prs: Presentación PowerPoint
            variable_text: Texto de la variable a buscar
            context_chars: Cantidad de caracteres de contexto antes/después
            
        Returns:
            Lista de diccionarios con 'before', 'variable', 'after', 'location'
        """
        contexts = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        full_text = ''.join(run.text for run in paragraph.runs)
                        contexts.extend(PatternDetector._extract_context_from_text(
                            full_text, variable_text, context_chars, 
                            f"Slide {slide_idx+1}, Forma {shape_idx+1}"
                        ))
        
        return contexts
    
    @staticmethod
    def _extract_context_from_text(full_text: str, variable_text: str, context_chars: int, location: str) -> List[Dict[str, str]]:
        """Extrae contexto de un texto específico"""
        contexts = []

        # Buscar todas las ocurrencias
        start_pos = 0
        while True:
            pos = full_text.find(variable_text, start_pos)
            if pos == -1:
                break

            # Extraer contexto
            before_start = max(0, pos - context_chars)
            before = full_text[before_start:pos]
            after_end = min(len(full_text), pos + len(variable_text) + context_chars)
            after = full_text[pos + len(variable_text):after_end]

            # Añadir elipsis si es necesario
            if before_start > 0:
                before = "..." + before
            if after_end < len(full_text):
                after = after + "..."

            contexts.append({
                'before': before,
                'variable': variable_text,
                'after': after,
                'location': location,
                'position': pos  # Guardar posición para reemplazo selectivo
            })

            start_pos = pos + 1

        return contexts

    @staticmethod
    def replace_text_by_context(text: str, search_text: str, replacement: str, context_indices: List[int] = None) -> str:
        """
        Reemplaza texto solo en los contextos especificados.

        Args:
            text: Texto completo
            search_text: Texto a buscar
            replacement: Texto de reemplazo
            context_indices: Lista de índices de ocurrencias a reemplazar (0-indexed).
                            Si es None, reemplaza todas.

        Returns:
            Texto modificado
        """
        if context_indices is None:
            # Reemplazar todas las ocurrencias
            return text.replace(search_text, replacement)

        # Encontrar todas las posiciones
        positions = []
        start_pos = 0
        while True:
            pos = text.find(search_text, start_pos)
            if pos == -1:
                break
            positions.append(pos)
            start_pos = pos + 1

        # Filtrar solo las posiciones indicadas
        positions_to_replace = [positions[i] for i in context_indices if i < len(positions)]

        # Reemplazar de atrás hacia adelante para mantener índices válidos
        for pos in reversed(positions_to_replace):
            text = text[:pos] + replacement + text[pos + len(search_text):]

        return text


class VariableNormalizer:
    """Normaliza nombres de variables"""
    
    @staticmethod
    def normalize_name(text: str) -> str:
        """Normaliza un texto a nombre de variable válido"""
        # Eliminar espacios y caracteres especiales
        normalized = re.sub(r'[^\w\s]', '', text)
        normalized = normalized.strip().lower()
        normalized = re.sub(r'\s+', '_', normalized)
        return normalized or "variable"
    
    @staticmethod
    def generate_default_question(var_name: str, var_type: str) -> str:
        """Genera una pregunta por defecto basada en el nombre y tipo de variable"""
        questions_map = {
            'texto': f"Ingrese el valor para {var_name.replace('_', ' ')}:",
            'numero': f"Ingrese el número para {var_name.replace('_', ' ')}:",
            'fecha': f"Ingrese la fecha para {var_name.replace('_', ' ')} (DD/MM/YYYY):",
            'hora': f"Ingrese la hora para {var_name.replace('_', ' ')} (HH:MM):",
            'lista': f"Seleccione una opción para {var_name.replace('_', ' ')}:",
            'email': f"Ingrese el email para {var_name.replace('_', ' ')}:",
            'telefono': f"Ingrese el teléfono para {var_name.replace('_', ' ')}:",
            'moneda': f"Ingrese el importe para {var_name.replace('_', ' ')}:",
            'porcentaje': f"Indique el porcentaje para {var_name.replace('_', ' ')}:",
            'booleano': f"Seleccione Sí/No para {var_name.replace('_', ' ')}:",
        }
        return questions_map.get(var_type, f"Ingrese el valor para {var_name.replace('_', ' ')}:")


class YAMLManager:
    """Gestiona archivos YAML de configuración"""
    
    @staticmethod
    def create_variable_config(variables: dict | List[Dict[str, Any]]) -> dict:
        """
        Crea la estructura YAML con todas las variables detectadas y sus metadatos.

        Ahora acepta listas de diccionarios (Fase 1) o el formato tradicional
        basado en diccionarios para mantener compatibilidad retro.
        Se añaden categorías, formatos sugeridos y banderas de calidad
        pensadas para los entregables diarios de firmas Big4.
        """

        # Compatibilidad: dict -> lista
        if isinstance(variables, dict):
            variable_items = []
            for var_id, info in variables.items():
                variable_items.append({"nombre": var_id, **info})
        else:
            variable_items = variables

        config = {
            'variables': [],
            'controles': {
                'ultima_generacion': datetime.utcnow().isoformat() + 'Z',
                'validador': 'plantillador_v2',
            },
        }

        for info in variable_items:
            if info.get('disabled', False):
                continue

            var_id = info.get('nombre')
            var_type = info.get('tipo', 'texto')
            pregunta = info.get('pregunta') or VariableNormalizer.generate_default_question(var_id, var_type)
            item = {
                'nombre': var_id,
                'tipo': var_type,
                'pregunta': pregunta,
                'categoria': info.get('categoria', 'general'),
                'requerido': info.get('requerido', True),
                'placeholder': info.get('placeholder'),
            }

            if info.get('opciones'):
                item['opciones'] = info['opciones']
            if info.get('meta'):
                item['meta'] = info['meta']
            if info.get('formato'):
                item['formato'] = info['formato']
            item['ejemplo'] = info.get('ejemplo') or YAMLManager._default_example(var_type)

            config['variables'].append(item)

        return config

    @staticmethod
    def _default_example(var_type: str) -> str:
        ejemplos = {
            'fecha': '31/12/2024',
            'hora': '14:30',
            'moneda': '15000.00',
            'porcentaje': '8.5%',
            'email': 'nombre@empresa.com',
            'telefono': '+34-600000000',
        }
        return ejemplos.get(var_type, 'valor')

    @staticmethod
    def create_value_manifest(values: Dict[str, str], schema: List[Dict[str, Any]], output_name: str) -> Dict[str, Any]:
        """Construye un manifiesto de reemplazo para auditoría."""

        indexed_schema = {item.get('nombre'): item for item in schema}
        applied = []
        for key, value in values.items():
            meta = indexed_schema.get(key, {})
            applied.append({
                'nombre': key,
                'tipo': meta.get('tipo', 'texto'),
                'categoria': meta.get('categoria', 'general'),
                'valor': value,
            })

        return {
            'informe': output_name,
            'generado_en': datetime.utcnow().isoformat() + 'Z',
            'valores': applied,
        }
    
    @staticmethod
    def save_yaml(data: Dict, filepath: str):
        """Guarda datos en archivo YAML"""
        with open(filepath, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
    
    @staticmethod
    def load_yaml(filepath: str) -> Dict:
        """Carga datos desde archivo YAML"""
        with open(filepath, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)


class DocumentProcessor:
    """Procesa documentos Word y PowerPoint"""
    
    @staticmethod
    def replace_in_docx(doc: Document, replacements: Dict[str, str]) -> Document:
        """Reemplaza variables en documento Word manteniendo formato"""
        
        # Reemplazar en párrafos
        for paragraph in doc.paragraphs:
            for var_id, value in replacements.items():
                placeholder = f"{{{{{var_id}}}}}"
                if placeholder in paragraph.text:
                    DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        # Reemplazar en tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for var_id, value in replacements.items():
                            placeholder = f"{{{{{var_id}}}}}"
                            if placeholder in paragraph.text:
                                DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        # Reemplazar en encabezados y pies
        for section in doc.sections:
            # Encabezado
            if section.header:
                for paragraph in section.header.paragraphs:
                    for var_id, value in replacements.items():
                        placeholder = f"{{{{{var_id}}}}}"
                        if placeholder in paragraph.text:
                            DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
            
            # Pie de página
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    for var_id, value in replacements.items():
                        placeholder = f"{{{{{var_id}}}}}"
                        if placeholder in paragraph.text:
                            DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        return doc
    
    @staticmethod
    def _replace_in_paragraph(paragraph, placeholder: str, value: str):
        """Reemplaza texto en un párrafo manteniendo el formato"""
        if placeholder not in paragraph.text:
            return
        
        # Estrategia: reconstruir el párrafo manteniendo los runs
        full_text = paragraph.text
        if placeholder in full_text:
            new_text = full_text.replace(placeholder, value)
            
            # Limpiar el párrafo
            for run in paragraph.runs:
                run.text = ''
            
            # Agregar el nuevo texto en el primer run
            if paragraph.runs:
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run(new_text)
    
    @staticmethod
    def replace_in_pptx(prs: Presentation, replacements: Dict[str, str]) -> Presentation:
        """Reemplaza variables en presentación PowerPoint usando el XML de los slides.

        Trabajar directamente sobre el XML evita errores en shapes que no tienen text_frame
        y permite reemplazar marcadores aunque estén divididos en varios runs.
        """

        namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }

        for slide in prs.slides:
            # Recorremos todos los párrafos del slide en el XML
            for paragraph in slide.part.element.xpath('.//a:p', namespaces=namespaces):
                text_elements = paragraph.xpath('.//a:t', namespaces=namespaces)
                if not text_elements:
                    continue

                full_text = ''.join(t.text or '' for t in text_elements)
                new_text = full_text

                for var_id, value in replacements.items():
                    placeholder = f"{{{{{var_id}}}}}"
                    if placeholder in new_text:
                        new_text = new_text.replace(placeholder, value)

                # Si hubo cambios, volcamos el texto nuevo en el primer nodo y vaciamos el resto
                if new_text != full_text:
                    text_elements[0].text = new_text
                    for extra in text_elements[1:]:
                        extra.text = ''

        return prs


class Validator:
    """Valida datos según el tipo"""
    
    @staticmethod
    def validate_email(email: str) -> bool:
        """Valida formato de email"""
        pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        return bool(re.match(pattern, email))
    
    @staticmethod
    def validate_phone(phone: str) -> bool:
        """Valida formato de teléfono (flexible)"""
        # Permite números, espacios, guiones y paréntesis
        pattern = r'^[\d\s\-\(\)\+]+$'
        return bool(re.match(pattern, phone)) and len(re.sub(r'\D', '', phone)) >= 9
    
    @staticmethod
    def validate_date(date_str: str) -> bool:
        """Valida formato de fecha DD/MM/YYYY"""
        pattern = r'^(0[1-9]|[12][0-9]|3[01])/(0[1-9]|1[0-2])/\d{4}$'
        return bool(re.match(pattern, date_str))
    
    @staticmethod
    def validate_time(time_str: str) -> bool:
        """Valida formato de hora HH:MM"""
        pattern = r'^([01][0-9]|2[0-3]):[0-5][0-9]$'
        return bool(re.match(pattern, time_str))
    
    @staticmethod
    def validate_number(number_str: str) -> bool:
        """Valida que sea un número"""
        try:
            float(number_str.replace(',', '.'))
            return True
        except ValueError:
            return False
