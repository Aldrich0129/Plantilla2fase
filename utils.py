"""
Librería común para el sistema de automatización de plantillas
Funciones compartidas entre Fase 1 y Fase 2
"""

import re
from typing import Dict, List, Set, Tuple, Any
from docx import Document
from docx.shared import RGBColor
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt
import yaml


class PatternDetector:
    """Detecta patrones de variables en documentos"""
    
    # Patrones de texto predefinidos
    TEXT_PATTERNS = {
        'llaves_simples': r'\{([^}]+)\}',
        'llaves_dobles': r'\{\{([^}]+)\}\}',
        'corchetes_simples': r'\[([^\]]+)\]',
        'corchetes_dobles': r'\[\[([^\]]+)\]\]',
    }
    
    @staticmethod
    def detect_colors_in_docx(doc: Document) -> Dict[str, Set[str]]:
        """Detecta todos los colores usados en un documento Word"""
        colors = {
            'text_colors': set(),
            'highlight_colors': set()
        }
        
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                # Color de texto
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    colors['text_colors'].add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                
                # Color de subrayado/resaltado
                if run.font.highlight_color:
                    colors['highlight_colors'].add(str(run.font.highlight_color))
        
        # También revisar tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run.font.color and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                colors['text_colors'].add(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
        
        return colors
    
    @staticmethod
    def detect_colors_in_pptx(prs: Presentation) -> Dict[str, Set[str]]:
        """
        Detecta todos los colores usados en una presentación PowerPoint.
        Maneja diferentes tipos de colores: RGB directo, colores de tema, etc.
        """
        colors = {
            'text_colors': set(),
            'highlight_colors': set()
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                # Verificar si el run tiene color
                                if not run.font.color:
                                    continue
                                
                                color_obj = run.font.color
                                rgb_value = None
                                
                                # Método 1: Intentar obtener RGB directo
                                if hasattr(color_obj, 'rgb'):
                                    try:
                                        rgb_value = color_obj.rgb
                                    except (AttributeError, TypeError):
                                        pass
                                
                                # Método 2: Si es color de tema, intentar obtener RGB del tema
                                if rgb_value is None and hasattr(color_obj, 'theme_color'):
                                    try:
                                        # Los colores de tema no se pueden convertir fácilmente
                                        # Se omiten por ahora
                                        pass
                                    except (AttributeError, TypeError):
                                        pass
                                
                                # Si obtuvimos un valor RGB válido, agregarlo
                                if rgb_value:
                                    colors['text_colors'].add(f"#{rgb_value[0]:02x}{rgb_value[1]:02x}{rgb_value[2]:02x}")
                                    
                            except Exception:
                                # Si hay cualquier otro error, simplemente continuar
                                # Esto evita que el programa se detenga por colores no soportados
                                pass
        
        return colors
    
    @staticmethod
    def extract_variables_by_pattern(text: str, pattern: str) -> List[str]:
        """Extrae variables según un patrón regex"""
        if pattern in PatternDetector.TEXT_PATTERNS:
            regex = PatternDetector.TEXT_PATTERNS[pattern]
        else:
            regex = pattern
        
        matches = re.findall(regex, text)
        return list(set(matches))  # Eliminar duplicados
    
    @staticmethod
    def extract_variables_by_color(doc: Document, color: str, color_type: str = 'text') -> List[Tuple[str, Any]]:
        """
        Extrae texto con un color específico, agrupando runs consecutivos del mismo color.
        
        Mejora clave: Si varios runs consecutivos tienen el mismo formato, 
        se agrupan en UNA sola variable.
        """
        variables = []
        
        # Procesar párrafos
        for paragraph in doc.paragraphs:
            grouped_vars = PatternDetector._group_consecutive_colored_runs(
                paragraph.runs, color, color_type
            )
            variables.extend(grouped_vars)
        
        # Procesar tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        grouped_vars = PatternDetector._group_consecutive_colored_runs(
                            paragraph.runs, color, color_type
                        )
                        variables.extend(grouped_vars)
        
        return variables
    
    @staticmethod
    def _group_consecutive_colored_runs(runs: list, color: str, color_type: str) -> List[Tuple[str, Any]]:
        """
        Agrupa runs consecutivos con el mismo color en una sola variable.
        
        Args:
            runs: Lista de runs del párrafo
            color: Color a buscar (hex para text, string para highlight)
            color_type: 'text' o 'highlight'
            
        Returns:
            Lista de tuplas (texto_agrupado, primer_run)
        """
        variables = []
        current_text = ""
        first_run = None
        
        for run in runs:
            run_matches = False
            
            # Verificar si el run tiene el color buscado
            if color_type == 'text' and run.font.color and run.font.color.rgb:
                rgb = run.font.color.rgb
                run_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                run_matches = (run_color == color)
            elif color_type == 'highlight' and run.font.highlight_color:
                run_matches = (str(run.font.highlight_color) == color)
            
            if run_matches:
                # Este run tiene el color buscado
                if first_run is None:
                    first_run = run
                current_text += run.text
            else:
                # Este run NO tiene el color, guardar lo acumulado
                if current_text.strip():
                    variables.append((current_text.strip(), first_run))
                current_text = ""
                first_run = None
        
        # Guardar el último grupo si existe
        if current_text.strip():
            variables.append((current_text.strip(), first_run))
        
        return variables
    
    @staticmethod
    def detect_list_options(text: str) -> List[str]:
        """
        Detecta si un texto contiene opciones separadas por delimitadores.
        Soporta: / | , y
        
        Ejemplos:
        "las cuentas anuales/los estados financieros/el balance" 
        → ["las cuentas anuales", "los estados financieros", "el balance"]
        """
        # Probar diferentes delimitadores
        for delimiter in ['/', '|']:
            if delimiter in text:
                options = [opt.strip() for opt in text.split(delimiter) if opt.strip()]
                # Solo considerar si hay al menos 2 opciones
                if len(options) >= 2:
                    return options
        
        return []


class VariableNormalizer:
    """Normaliza nombres de variables"""
    
    @staticmethod
    def normalize_name(text: str) -> str:
        """Convierte texto en un nombre de variable válido"""
        # Eliminar espacios y caracteres especiales
        normalized = re.sub(r'[^\w\s]', '', text)
        normalized = normalized.strip().lower()
        normalized = re.sub(r'\s+', '_', normalized)
        return normalized or "variable"
    
    @staticmethod
    def generate_default_question(var_name: str, var_type: str) -> str:
        """Genera una pregunta por defecto basada en el nombre y tipo de variable"""
        questions_map = {
            'texto': f"Ingrese el valor para {var_name.replace('_', ' ')}:",
            'numero': f"Ingrese el número para {var_name.replace('_', ' ')}:",
            'fecha': f"Ingrese la fecha para {var_name.replace('_', ' ')} (DD/MM/YYYY):",
            'hora': f"Ingrese la hora para {var_name.replace('_', ' ')} (HH:MM):",
            'lista': f"Seleccione una opción para {var_name.replace('_', ' ')}:",
            'email': f"Ingrese el email para {var_name.replace('_', ' ')}:",
            'telefono': f"Ingrese el teléfono para {var_name.replace('_', ' ')}:",
            'moneda': f"Ingrese el importe para {var_name.replace('_', ' ')}:"
        }
        return questions_map.get(var_type, f"Ingrese el valor para {var_name.replace('_', ' ')}:")


class YAMLManager:
    """Gestiona archivos YAML de configuración"""
    
    @staticmethod
    def create_variable_config(variables: dict) -> dict:
        """
        Crea la estructura YAML con todas las variables detectadas y sus metadatos.
        Compatible con tipos especiales: lista, moneda, telefono.
        """
        config = {'variables': []}

        for var_id, info in variables.items():
            item = {
                'nombre': var_id,
                'tipo': info.get('tipo', 'texto'),
                'pregunta': (
                    info.get('pregunta') or
                    VariableNormalizer.generate_default_question(var_id, info.get('tipo', 'texto'))
                )
            }

            # Si es lista → añadir opciones
            if info.get('tipo') == 'lista' and info.get('opciones'):
                item['opciones'] = info['opciones']

            # 🔹 Si tiene metadatos (moneda o teléfono), los guardamos en YAML
            if info.get('meta'):
                item['meta'] = info['meta']

            config['variables'].append(item)

        return config
    
    @staticmethod
    def save_yaml(data: Dict, filepath: str):
        """Guarda datos en archivo YAML"""
        with open(filepath, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
    
    @staticmethod
    def load_yaml(filepath: str) -> Dict:
        """Carga datos desde archivo YAML"""
        with open(filepath, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)


class DocumentProcessor:
    """Procesa documentos Word y PowerPoint"""
    
    @staticmethod
    def replace_in_docx(doc: Document, replacements: Dict[str, str]) -> Document:
        """Reemplaza variables en documento Word manteniendo formato"""
        
        # Reemplazar en párrafos
        for paragraph in doc.paragraphs:
            for var_id, value in replacements.items():
                placeholder = f"{{{{{var_id}}}}}"
                if placeholder in paragraph.text:
                    DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        # Reemplazar en tablas
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for var_id, value in replacements.items():
                            placeholder = f"{{{{{var_id}}}}}"
                            if placeholder in paragraph.text:
                                DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        # Reemplazar en encabezados y pies
        for section in doc.sections:
            # Encabezado
            if section.header:
                for paragraph in section.header.paragraphs:
                    for var_id, value in replacements.items():
                        placeholder = f"{{{{{var_id}}}}}"
                        if placeholder in paragraph.text:
                            DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
            
            # Pie de página
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    for var_id, value in replacements.items():
                        placeholder = f"{{{{{var_id}}}}}"
                        if placeholder in paragraph.text:
                            DocumentProcessor._replace_in_paragraph(paragraph, placeholder, value)
        
        return doc
    
    @staticmethod
    def _replace_in_paragraph(paragraph, placeholder: str, value: str):
        """Reemplaza texto en un párrafo manteniendo el formato"""
        if placeholder not in paragraph.text:
            return
        
        # Estrategia: reconstruir el párrafo manteniendo los runs
        full_text = paragraph.text
        if placeholder in full_text:
            new_text = full_text.replace(placeholder, value)
            
            # Limpiar el párrafo
            for run in paragraph.runs:
                run.text = ''
            
            # Agregar el nuevo texto en el primer run
            if paragraph.runs:
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run(new_text)
    
    @staticmethod
    def replace_in_pptx(prs: Presentation, replacements: Dict[str, str]) -> Presentation:
        """Reemplaza variables en presentación PowerPoint manteniendo formato"""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    DocumentProcessor._replace_in_textframe(shape.text_frame, replacements)
                
                # También revisar tablas
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            DocumentProcessor._replace_in_textframe(cell.text_frame, replacements)
        
        return prs
    
    @staticmethod
    def _replace_in_textframe(text_frame, replacements: Dict[str, str]):
        """Reemplaza texto en un text_frame de PowerPoint"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                for var_id, value in replacements.items():
                    placeholder = f"{{{{{var_id}}}}}"
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)


class Validator:
    """Valida datos según el tipo"""
    
    @staticmethod
    def validate_email(email: str) -> bool:
        """Valida formato de email"""
        pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        return bool(re.match(pattern, email))
    
    @staticmethod
    def validate_phone(phone: str) -> bool:
        """Valida formato de teléfono (flexible)"""
        # Permite números, espacios, guiones y paréntesis
        pattern = r'^[\d\s\-\(\)\+]+$'
        return bool(re.match(pattern, phone)) and len(re.sub(r'\D', '', phone)) >= 9
    
    @staticmethod
    def validate_date(date_str: str) -> bool:
        """Valida formato de fecha DD/MM/YYYY"""
        pattern = r'^(0[1-9]|[12][0-9]|3[01])/(0[1-9]|1[0-2])/\d{4}$'
        return bool(re.match(pattern, date_str))
    
    @staticmethod
    def validate_time(time_str: str) -> bool:
        """Valida formato de hora HH:MM"""
        pattern = r'^([01][0-9]|2[0-3]):[0-5][0-9]$'
        return bool(re.match(pattern, time_str))
    
    @staticmethod
    def validate_number(number_str: str) -> bool:
        """Valida que sea un número"""
        try:
            float(number_str.replace(',', '.'))
            return True
        except ValueError:
            return False