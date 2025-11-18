"""Lógica de la Fase 1 desacoplada de la UI."""
from __future__ import annotations

import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

from docx import Document
from pptx import Presentation

from utils_v2 import PatternDetector, VariableNormalizer, YAMLManager
from .text_utils import clean_pattern_markers, infer_variable_type, sanitize_placeholders


@dataclass
class DetectedVariable:
    nombre: str
    marcador_original: str
    tipo: str


class TemplateBuilder:
    """Procesa un documento y genera la plantilla normalizada y el YAML."""

    def __init__(self, workdir: Path):
        self.workdir = Path(workdir)
        self.workdir.mkdir(parents=True, exist_ok=True)

    def save_upload(self, file_bytes, filename: str) -> Path:
        target = self.workdir / filename
        with open(target, "wb") as fh:
            fh.write(file_bytes)
        return target

    def detect_variables(self, file_path: Path) -> Tuple[List[DetectedVariable], Dict[str, str]]:
        content_blocks = list(self._iter_text_blocks(file_path))
        replacements: Dict[str, str] = {}
        detected: Dict[str, DetectedVariable] = {}

        for block in content_blocks:
            sanitized = sanitize_placeholders(block)
            for pattern in PatternDetector.TEXT_PATTERNS.values():
                for match in PatternDetector.extract_variables_by_pattern(sanitized, pattern):
                    cleaned = clean_pattern_markers(match)
                    normalized = VariableNormalizer.normalize_name(cleaned)
                    placeholder_original = self._wrap_placeholder(match)
                    placeholder_normalized = f"{{{{{normalized}}}}}"
                    replacements[placeholder_original] = placeholder_normalized
                    if normalized not in detected:
                        detected[normalized] = DetectedVariable(
                            nombre=normalized,
                            marcador_original=placeholder_original,
                            tipo=infer_variable_type(cleaned),
                        )

        return list(detected.values()), replacements

    def build_outputs(self, file_path: Path) -> Tuple[Path, Path, List[DetectedVariable]]:
        detected_vars, replacements = self.detect_variables(file_path)
        if file_path.suffix.lower() == ".docx":
            template_path = self._generate_docx_template(file_path, replacements)
        elif file_path.suffix.lower() == ".pptx":
            template_path = self._generate_pptx_template(file_path, replacements)
        else:
            raise ValueError("Formato no soportado. Usa .docx o .pptx")

        yaml_path = self.workdir / f"{file_path.stem}_variables.yaml"
        yaml_config = YAMLManager.create_variable_config({
            var.nombre: {"tipo": var.tipo} for var in detected_vars
        })
        YAMLManager.save_yaml(yaml_config, str(yaml_path))
        return template_path, yaml_path, detected_vars

    # Helpers

    def _iter_text_blocks(self, file_path: Path) -> Iterable[str]:
        suffix = file_path.suffix.lower()
        if suffix == ".docx":
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                yield paragraph.text
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            yield paragraph.text
        elif suffix == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        yield shape.text
        else:
            raise ValueError("Formato no soportado")

    @staticmethod
    def _wrap_placeholder(content: str) -> str:
        # Mantener dobles llaves si ya estaban presentes
        if content.startswith("{{") and content.endswith("}}"):
            return content
        if content.startswith("{") and content.endswith("}"):
            return content
        if content.startswith("[") and content.endswith("]"):
            return content
        return f"{{{{{clean_pattern_markers(content)}}}}}"

    def _generate_docx_template(self, file_path: Path, replacements: Dict[str, str]) -> Path:
        doc = Document(file_path)
        self._replace_text_in_docx(doc, replacements)
        target = self.workdir / f"{file_path.stem}_template.docx"
        doc.save(target)
        return target

    def _replace_text_in_docx(self, doc: Document, replacements: Dict[str, str]) -> None:
        def replace_in_paragraph(paragraph, needle: str, replacement: str):
            text = paragraph.text
            if needle not in text:
                return
            new_text = text.replace(needle, replacement)
            if paragraph.runs:
                for run in paragraph.runs:
                    run.text = ""
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run(new_text)

        for paragraph in doc.paragraphs:
            for src, tgt in replacements.items():
                replace_in_paragraph(paragraph, src, tgt)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for src, tgt in replacements.items():
                            replace_in_paragraph(paragraph, src, tgt)

        for section in doc.sections:
            if section.header:
                for paragraph in section.header.paragraphs:
                    for src, tgt in replacements.items():
                        replace_in_paragraph(paragraph, src, tgt)
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    for src, tgt in replacements.items():
                        replace_in_paragraph(paragraph, src, tgt)

    def _generate_pptx_template(self, file_path: Path, replacements: Dict[str, str]) -> Path:
        prs = Presentation(file_path)
        namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for slide in prs.slides:
            for paragraph in slide.part.element.xpath('.//a:p', namespaces=namespaces):
                text_elements = paragraph.xpath('.//a:t', namespaces=namespaces)
                if not text_elements:
                    continue
                full_text = "".join(t.text or "" for t in text_elements)
                new_text = full_text
                for src, tgt in replacements.items():
                    if src in new_text:
                        new_text = new_text.replace(src, tgt)
                if new_text != full_text:
                    text_elements[0].text = new_text
                    for extra in text_elements[1:]:
                        extra.text = ""
        target = self.workdir / f"{file_path.stem}_template.pptx"
        prs.save(target)
        return target


__all__ = ["TemplateBuilder", "DetectedVariable"]
