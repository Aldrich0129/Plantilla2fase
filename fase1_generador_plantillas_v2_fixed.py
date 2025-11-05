"""
FASE 1: Generador de Plantillas v2.0 FINAL
Detecta variables en documentos Word/PPT y genera plantillas normalizadas

NUEVAS FUNCIONALIDADES:
- Patrón combinado (AND entre dos patrones)
- Desidentificar variables
- Mostrar contexto de variables
- División libre con selector manual
- Dividir variables por contexto seleccionado
- Detección mejorada de fechas con "de" (día de mes de año, etc.)
"""

import streamlit as st
import os
import tempfile
import re
from pathlib import Path
from docx import Document
from pptx import Presentation
from utils_v2 import (
    PatternDetector, VariableNormalizer, YAMLManager, DocumentProcessor
)
import io
import unicodedata

_SANITIZE_PATTERNS = [
    (r'[\[\(\{]\s*\{\{\s*([^\}]+?)\s*\}\}\s*[\]\)\}]', r'{{\1}}'),
    (r'\{\{\s*([^\}]+?)\s*\}\}\s*[\]\}\)]', r'{{\1}}'),
    (r'[\[\(\{]\s*\{\{\s*([^\}]+?)\s*\}\}', r'{{\1}}'),
    (r'\{\{\s*([^\}]+?)\s*\}\}', lambda m: '{{' + re.sub(r'\s+', '_', m.group(1).strip()) + '}}'),
]

def sanitize_placeholders(text: str) -> str:
    """Limpia artefactos como '{{var}}]' o '{{var}}}' y normaliza '{{ var }}' -> '{{var}}'."""
    t = unicodedata.normalize('NFKC', text).replace('\u00A0', ' ')
    changed = True
    while changed:
        changed = False
        for pat, repl in _SANITIZE_PATTERNS:
            new_t = re.sub(pat, repl, t)
            if new_t != t:
                t = new_t
                changed = True
    return t

st.set_page_config(
    page_title="Fase 1: Generador de Plantillas v2.0",
    page_icon="📝",
    layout="wide"
)

# Inicializar estado de sesión
if 'variables' not in st.session_state:
    st.session_state.variables = {}
if 'detected_patterns' not in st.session_state:
    st.session_state.detected_patterns = {}
if 'doc_type' not in st.session_state:
    st.session_state.doc_type = None
if 'original_doc' not in st.session_state:
    st.session_state.original_doc = None
if 'work_dir' not in st.session_state:
    st.session_state.work_dir = tempfile.mkdtemp()
if 'selected_for_merge' not in st.session_state:
    st.session_state.selected_for_merge = []


def hex_to_color_name(hex_color):
    """Convierte código hex a nombre de color en español"""
    color_map = {
        '#ff0000': '🔴 Rojo',
        '#00ff00': '🟢 Verde Lima',
        '#0000ff': '🔵 Azul',
        '#ffff00': '🟡 Amarillo',
        '#ff00ff': '🟣 Magenta',
        '#00ffff': '🔵 Cian',
        '#000000': '⚫ Negro',
        '#ffffff': '⚪ Blanco',
        '#808080': '⚪ Gris',
        '#800000': '🔴 Rojo Oscuro',
        '#008000': '🟢 Verde',
        '#000080': '🔵 Azul Marino',
        '#808000': '🟡 Oliva',
        '#800080': '🟣 Púrpura',
        '#008080': '🔵 Verde Azulado',
        '#c0c0c0': '⚪ Plata',
        '#ff6600': '🟠 Naranja',
        '#ff9900': '🟠 Naranja Claro',
        '#993300': '🟤 Marrón',
        '#660000': '🔴 Granate',
        '#006600': '🟢 Verde Oscuro',
        '#003366': '🔵 Azul Oscuro',
        '#663399': '🟣 Púrpura Medio',
        '#336699': '🔵 Azul Acero',
    }
    
    if hex_color.lower() in color_map:
        return color_map[hex_color.lower()]
    
    hex_lower = hex_color.lower().replace('#', '')
    if len(hex_lower) == 6:
        r = int(hex_lower[0:2], 16)
        g = int(hex_lower[2:4], 16)
        b = int(hex_lower[4:6], 16)
        
        if r > 200 and g < 100 and b < 100:
            return '🔴 Rojo'
        elif r < 100 and g > 200 and b < 100:
            return '🟢 Verde'
        elif r < 100 and g < 100 and b > 200:
            return '🔵 Azul'
        elif r > 200 and g > 200 and b < 100:
            return '🟡 Amarillo'
        elif r > 200 and g < 100 and b > 200:
            return '🟣 Magenta'
        elif r < 100 and g > 200 and b > 200:
            return '🔵 Cian'
        elif r > 150 and g > 100 and b < 100:
            return '🟠 Naranja'
        elif r > 100 and g < 100 and b > 150:
            return '🟣 Púrpura'
        elif r < 50 and g < 50 and b < 50:
            return '⚫ Negro'
        elif r > 200 and g > 200 and b > 200:
            return '⚪ Blanco'
        else:
            return '⚪ Gris'
    
    return hex_color


def clean_pattern_markers(text: str) -> str:
    """Elimina todos los marcadores de patrones ({}, [], etc.) del texto"""
    text = re.sub(r'\{\{([^}]+)\}\}', r'\1', text)
    text = re.sub(r'\{([^}]+)\}', r'\1', text)
    text = re.sub(r'\[\[([^\]]+)\]\]', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]', r'\1', text)
    return text.strip()


def infer_variable_type(text: str) -> str:
    text_lower = text.lower()

    if any(tok in text_lower for tok in ['€', 'eur', 'euro', 'euros', 'usd', '$', 'dólar', 'dolar', 'dólares', 'dolares', 'importe', 'monto']):
        return 'moneda'
    if any(word in text_lower for word in ['fecha', 'date', 'día', 'mes', 'año']):
        return 'fecha'
    elif any(word in text_lower for word in ['hora', 'time', 'horario']):
        return 'hora'
    elif any(word in text_lower for word in ['email', 'correo', 'mail']):
        return 'email'
    elif any(word in text_lower for word in ['teléfono', 'telefono', 'phone', 'móvil', 'movil']):
        return 'telefono'
    elif any(word in text_lower for word in ['número', 'numero', 'cantidad', 'monto', 'precio']):
        return 'numero'
    else:
        return 'texto'


def merge_variables(var_ids: list, new_name: str = None):
    """Fusiona múltiples variables en una sola"""
    if len(var_ids) < 2:
        st.error("Debes seleccionar al menos 2 variables para fusionar")
        return
    
    merged_text = " ".join([
        st.session_state.variables[vid]['original_text'] 
        for vid in var_ids if vid in st.session_state.variables
    ])
    
    if not new_name:
        new_name = "_".join(var_ids[:2])
    
    first_var = st.session_state.variables[var_ids[0]]
    
    st.session_state.variables[new_name] = {
        'original_text': merged_text,
        'tipo': first_var['tipo'],
        'pattern': first_var['pattern'],
        'pregunta': first_var.get('pregunta', ''),
        'opciones': first_var.get('opciones', []),
        'disabled': False
    }
    
    for vid in var_ids:
        if vid in st.session_state.variables:
            del st.session_state.variables[vid]
    
    st.session_state.selected_for_merge = []
    st.success(f"✅ Variables fusionadas en: `{new_name}`")
    st.rerun()


def split_variable(var_id: str, delimiter: str = '/'):
    """Divide una variable en múltiples variables según un delimitador"""
    if var_id not in st.session_state.variables:
        return
    
    var_info = st.session_state.variables[var_id]
    original_text = var_info['original_text']
    
    parts = [p.strip() for p in original_text.split(delimiter) if p.strip()]
    
    if len(parts) < 2:
        st.error(f"No se encontraron múltiples partes separadas por '{delimiter}'")
        return
    
    del st.session_state.variables[var_id]
    
    for i, part in enumerate(parts):
        new_var_id = VariableNormalizer.normalize_name(part)
        if new_var_id in st.session_state.variables:
            new_var_id = f"{new_var_id}_{i+1}"
        
        st.session_state.variables[new_var_id] = {
            'original_text': part,
            'tipo': var_info['tipo'],
            'pattern': var_info['pattern'],
            'pregunta': '',
            'opciones': [],
            'disabled': False
        }
    
    st.success(f"✅ Variable dividida en {len(parts)} nuevas variables")
    st.rerun()


def split_variable_free(var_id: str, start_idx: int, end_idx: int, new_var_name: str):
    """Divide una variable usando selección libre por índices"""
    if var_id not in st.session_state.variables:
        return

    var_info = st.session_state.variables[var_id]
    original_text = var_info['original_text']

    if start_idx < 0 or end_idx > len(original_text) or start_idx >= end_idx:
        st.error("Índices inválidos. Verifica los valores ingresados.")
        return

    selected_part = original_text[start_idx:end_idx]
    remaining_part = original_text[:start_idx] + original_text[end_idx:]

    if not selected_part.strip() or not remaining_part.strip():
        st.error("Una de las partes resultó vacía. Ajusta los índices.")
        return

    new_var_id = VariableNormalizer.normalize_name(new_var_name if new_var_name else selected_part)

    suffix = 1
    original_new_var_id = new_var_id
    while new_var_id in st.session_state.variables:
        new_var_id = f"{original_new_var_id}_{suffix}"
        suffix += 1

    st.session_state.variables[new_var_id] = {
        'original_text': selected_part.strip(),
        'tipo': var_info['tipo'],
        'pattern': var_info['pattern'],
        'pregunta': '',
        'opciones': [],
        'disabled': False
    }

    st.session_state.variables[var_id]['original_text'] = remaining_part.strip()

    st.success(f"✅ Variable dividida en: `{new_var_id}` y `{var_id}` (modificada)")
    st.rerun()


def split_variable_by_context(var_id: str, contexts_groups: dict):
    """
    Divide una variable en múltiples variables según los contextos seleccionados.

    Args:
        var_id: ID de la variable a dividir
        contexts_groups: Dict donde cada key es el nuevo nombre de variable y
                        el value es la lista de índices de contextos
    Ejemplo: {'mes_de_cierre': [0, 1, 2], 'mes_de_facturacion': [3, 4]}
    """
    if var_id not in st.session_state.variables:
        st.error(f"Variable {var_id} no encontrada")
        return

    if not contexts_groups or len(contexts_groups) < 2:
        st.error("Debes crear al menos 2 grupos de contextos")
        return

    var_info = st.session_state.variables[var_id]

    # Crear las nuevas variables para cada grupo
    for new_var_name, context_indices in contexts_groups.items():
        if not context_indices:
            continue

        new_var_id = VariableNormalizer.normalize_name(new_var_name)

        # Evitar duplicados
        suffix = 1
        original_new_var_id = new_var_id
        while new_var_id in st.session_state.variables:
            new_var_id = f"{original_new_var_id}_{suffix}"
            suffix += 1

        # Crear la nueva variable con los contextos específicos
        st.session_state.variables[new_var_id] = {
            'original_text': var_info['original_text'],
            'tipo': var_info['tipo'],
            'pattern': var_info['pattern'],
            'pregunta': '',
            'opciones': var_info.get('opciones', []),
            'disabled': False,
            'context_indices': context_indices  # Guardar los índices de contexto
        }

    # Eliminar la variable original
    del st.session_state.variables[var_id]

    st.success(f"✅ Variable `{var_id}` dividida en {len(contexts_groups)} variables por contexto")
    st.rerun()


def toggle_variable_enabled(var_id: str):
    """Habilita o deshabilita una variable"""
    if var_id in st.session_state.variables:
        current_state = st.session_state.variables[var_id].get('disabled', False)
        st.session_state.variables[var_id]['disabled'] = not current_state
        st.rerun()


def detect_combined_pattern_variables(doc, file_extension: str, pattern1_type: str, pattern1_value: str, 
                                       pattern2_type: str, pattern2_value: str):
    """Detecta variables que cumplan AMBOS patrones (AND)"""
    variables_found = {}
    
    if file_extension == 'docx':
        full_text = extract_text_from_docx(doc)
    else:
        full_text = extract_text_from_pptx(doc)
    
    # Detectar variables del patrón 1
    vars_pattern1 = set()
    if pattern1_type == 'text':
        vars_list = PatternDetector.extract_variables_by_pattern(full_text, pattern1_value)
        vars_pattern1 = set(vars_list)
    elif pattern1_type == 'text_color' and file_extension == 'docx':
        color_vars = PatternDetector.extract_variables_by_color(doc, pattern1_value, 'text')
        vars_pattern1 = {var_text for var_text, _ in color_vars}
    elif pattern1_type == 'highlight_color' and file_extension == 'docx':
        color_vars = PatternDetector.extract_variables_by_color(doc, pattern1_value, 'highlight')
        vars_pattern1 = {var_text for var_text, _ in color_vars}
    
    # Detectar variables del patrón 2
    vars_pattern2 = set()
    if pattern2_type == 'text':
        vars_list = PatternDetector.extract_variables_by_pattern(full_text, pattern2_value)
        vars_pattern2 = set(vars_list)
    elif pattern2_type == 'text_color' and file_extension == 'docx':
        color_vars = PatternDetector.extract_variables_by_color(doc, pattern2_value, 'text')
        vars_pattern2 = {var_text for var_text, _ in color_vars}
    elif pattern2_type == 'highlight_color' and file_extension == 'docx':
        color_vars = PatternDetector.extract_variables_by_color(doc, pattern2_value, 'highlight')
        vars_pattern2 = {var_text for var_text, _ in color_vars}
    
    # Intersección: solo variables en AMBOS conjuntos
    common_vars = vars_pattern1.intersection(vars_pattern2)
    
    for var_text in common_vars:
        var_id = VariableNormalizer.normalize_name(var_text)
        if var_id not in variables_found:
            detected_options = PatternDetector.detect_list_options(var_text)
            var_tipo = 'lista' if detected_options else infer_variable_type(var_text)
            
            variables_found[var_id] = {
                'original_text': var_text,
                'tipo': var_tipo,
                'pattern': f'combined_{pattern1_type}_{pattern2_type}',
                'pregunta': '',
                'opciones': detected_options if detected_options else [],
                'disabled': False
            }
    
    return variables_found


def main():
    st.title("📝 Fase 1: Generador de Plantillas v2.0")
    st.markdown("### 🆕 Nuevas funcionalidades:")
    st.info("🔗 Patrón combinado | 🗑️ Desactivar variables | 📍 Contexto | ✂️ División libre | 🎯 División por contexto | 📅 Fechas con 'de'")
    st.markdown("---")

    with st.sidebar:
        st.header("📖 Guía de Uso")
        st.markdown("""
        ### Pasos:
        1. **Subir** documento
        2. **Seleccionar** patrones
        3. **(Opcional)** Patrón combinado
        4. **Revisar** variables
        5. **Configurar** y exportar

        ### 🆕 Nuevas funciones:
        - 🔗 Patrón combinado (AND)
        - 🗑️ Desactivar variables
        - 📍 Ver contexto
        - ✂️ División libre
        - 🎯 División por contexto
        - 📅 Fechas con "de"
        """)
    
    # Paso 1: Subir documento
    st.header("1️⃣ Subir Documento")
    uploaded_file = st.file_uploader(
        "Seleccione un archivo Word (.docx) o PowerPoint (.pptx)",
        type=['docx', 'pptx']
    )
    
    if uploaded_file is not None:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        st.session_state.doc_type = file_extension
        
        file_bytes = uploaded_file.read()
        
        if file_extension == 'docx':
            doc = Document(io.BytesIO(file_bytes))
            st.session_state.original_doc = doc
            st.success(f"✅ Documento Word cargado: {uploaded_file.name}")
        else:
            prs = Presentation(io.BytesIO(file_bytes))
            st.session_state.original_doc = prs
            st.success(f"✅ Presentación PowerPoint cargada: {uploaded_file.name}")
        
        st.markdown("---")
        
        # Paso 2: Detectar patrones (DISEÑO ORIGINAL)
        st.header("2️⃣ Detectar Patrones")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("Patrones de Texto")
            text_patterns = st.multiselect(
                "Seleccione patrones de texto:",
                options=list(PatternDetector.TEXT_PATTERNS.keys()),
                default=[],
                format_func=lambda x: {
                    'llaves_simples': '{variable}',
                    'llaves_dobles': '{{variable}}',
                    'corchetes_simples': '[variable]',
                    'corchetes_dobles': '[[variable]]'
                }.get(x, x)
            )
        
        with col2:
            st.subheader("Patrones de Color")
            
            if file_extension == 'docx':
                colors = PatternDetector.detect_colors_in_docx(doc)
            else:
                colors = PatternDetector.detect_colors_in_pptx(prs)
            
            text_colors_list = sorted(list(colors['text_colors']))
            highlight_colors_list = sorted(list(colors['highlight_colors']))
            
            selected_text_colors = []
            selected_highlight_colors = []
            
            if text_colors_list:
                st.write("**🎨 Colores de texto:**")
                for color in text_colors_list:
                    color_name = hex_to_color_name(color)
                    col_a, col_b = st.columns([0.15, 0.85])
                    with col_a:
                        st.markdown(
                            f'<div style="width:30px; height:30px; background-color:{color}; '
                            f'border: 2px solid #ccc; border-radius: 4px;"></div>',
                            unsafe_allow_html=True
                        )
                    with col_b:
                        if st.checkbox(f"{color_name}", key=f"text_{color}"):
                            selected_text_colors.append(color)
            
            if highlight_colors_list and file_extension == 'docx':
                st.write("**✏️ Colores de subrayado:**")
                for color in highlight_colors_list:
                    if st.checkbox(f"Subrayado: {color}", key=f"highlight_{color}"):
                        selected_highlight_colors.append(color)
            
            if not text_colors_list and not highlight_colors_list:
                st.info("ℹ️ No se detectaron colores especiales")
        
        # 🆕 Patrón combinado (EN EXPANDER SEPARADO)
        st.markdown("---")
        with st.expander("🔗 Crear Patrón Combinado (AND)", expanded=False):
            st.markdown("**Detecta variables que cumplan AMBOS patrones simultáneamente**")
            
            use_combined = st.checkbox("✅ Activar patrón combinado", key="use_combined")
            
            if use_combined:
                comb_col1, comb_col2 = st.columns(2)
                
                with comb_col1:
                    st.markdown("**Patrón 1:**")
                    pattern1_type = st.radio(
                        "Tipo:",
                        ['text', 'text_color', 'highlight_color'],
                        format_func=lambda x: {'text': '📝 Formato', 'text_color': '🎨 Color texto', 'highlight_color': '✏️ Subrayado'}.get(x),
                        key="p1_type"
                    )
                    
                    if pattern1_type == 'text':
                        pattern1_value = st.selectbox(
                            "Formato:",
                            list(PatternDetector.TEXT_PATTERNS.keys()),
                            format_func=lambda x: {
                                'llaves_simples': '{variable}',
                                'llaves_dobles': '{{variable}}',
                                'corchetes_simples': '[variable]',
                                'corchetes_dobles': '[[variable]]'
                            }.get(x),
                            key="p1_val"
                        )
                    elif pattern1_type == 'text_color':
                        pattern1_value = st.selectbox("Color:", text_colors_list, format_func=hex_to_color_name, key="p1_col") if text_colors_list else None
                    else:
                        pattern1_value = st.selectbox("Subrayado:", highlight_colors_list, key="p1_high") if highlight_colors_list else None
                
                with comb_col2:
                    st.markdown("**Patrón 2:**")
                    pattern2_type = st.radio(
                        "Tipo:",
                        ['text', 'text_color', 'highlight_color'],
                        format_func=lambda x: {'text': '📝 Formato', 'text_color': '🎨 Color texto', 'highlight_color': '✏️ Subrayado'}.get(x),
                        key="p2_type"
                    )
                    
                    if pattern2_type == 'text':
                        pattern2_value = st.selectbox(
                            "Formato:",
                            list(PatternDetector.TEXT_PATTERNS.keys()),
                            format_func=lambda x: {
                                'llaves_simples': '{variable}',
                                'llaves_dobles': '{{variable}}',
                                'corchetes_simples': '[variable]',
                                'corchetes_dobles': '[[variable]]'
                            }.get(x),
                            key="p2_val"
                        )
                    elif pattern2_type == 'text_color':
                        pattern2_value = st.selectbox("Color:", text_colors_list, format_func=hex_to_color_name, key="p2_col") if text_colors_list else None
                    else:
                        pattern2_value = st.selectbox("Subrayado:", highlight_colors_list, key="p2_high") if highlight_colors_list else None
                
                if pattern1_value and pattern2_value:
                    st.success(f"✅ Configurado: **{pattern1_type}** AND **{pattern2_type}**")
        
        # Botón detectar (LÓGICA ORIGINAL DEL CÓDIGO QUE FUNCIONA)
        st.markdown("---")
        if st.button("🔍 Detectar Variables", type="primary"):
            variables_found = {}

            # Extraer texto completo
            if file_extension == 'docx':
                full_text = extract_text_from_docx(doc)
            else:
                full_text = extract_text_from_pptx(doc)

            # 🆕 Detectar patrones de fecha con "de" ANTES de los patrones normales
            date_patterns = PatternDetector.detect_date_with_de_patterns(full_text)
            for date_text in date_patterns:
                var_id = VariableNormalizer.normalize_name(date_text)
                if var_id not in variables_found:
                    variables_found[var_id] = {
                        'original_text': date_text,
                        'tipo': 'fecha',
                        'pattern': 'date_with_de',
                        'pregunta': '',
                        'opciones': [],
                        'disabled': False
                    }

            # Detectar por patrones de texto (IGUAL QUE ORIGINAL)
            for pattern in text_patterns:
                vars_list = PatternDetector.extract_variables_by_pattern(full_text, pattern)
                for var_text in vars_list:
                    var_id = VariableNormalizer.normalize_name(var_text)
                    if var_id not in variables_found:
                        detected_options = PatternDetector.detect_list_options(var_text)
                        var_tipo = 'lista' if detected_options else infer_variable_type(var_text)

                        variables_found[var_id] = {
                            'original_text': var_text,
                            'tipo': var_tipo,
                            'pattern': pattern,
                            'pregunta': '',
                            'opciones': detected_options if detected_options else [],
                            'disabled': False
                        }
            
            # Detectar por colores (IGUAL QUE ORIGINAL)
            if file_extension == 'docx':
                for color in selected_text_colors:
                    color_vars = PatternDetector.extract_variables_by_color(doc, color, 'text')
                    for var_text, run in color_vars:
                        var_id = VariableNormalizer.normalize_name(var_text)
                        if var_id not in variables_found:
                            detected_options = PatternDetector.detect_list_options(var_text)
                            var_tipo = 'lista' if detected_options else infer_variable_type(var_text)
                            
                            variables_found[var_id] = {
                                'original_text': var_text,
                                'tipo': var_tipo,
                                'pattern': f'color_texto_{color}',
                                'pregunta': '',
                                'opciones': detected_options if detected_options else [],
                                'disabled': False
                            }
                
                for color in selected_highlight_colors:
                    color_vars = PatternDetector.extract_variables_by_color(doc, color, 'highlight')
                    for var_text, run in color_vars:
                        var_id = VariableNormalizer.normalize_name(var_text)
                        if var_id not in variables_found:
                            detected_options = PatternDetector.detect_list_options(var_text)
                            var_tipo = 'lista' if detected_options else infer_variable_type(var_text)
                            
                            variables_found[var_id] = {
                                'original_text': var_text,
                                'tipo': var_tipo,
                                'pattern': f'color_subrayado_{color}',
                                'pregunta': '',
                                'opciones': detected_options if detected_options else [],
                                'disabled': False
                            }
            
            # 🆕 Detectar con patrón combinado si está activado
            if use_combined and pattern1_value and pattern2_value:
                combined_vars = detect_combined_pattern_variables(
                    doc if file_extension == 'docx' else prs,
                    file_extension,
                    pattern1_type, pattern1_value,
                    pattern2_type, pattern2_value
                )
                
                for var_id, var_info in combined_vars.items():
                    if var_id in variables_found:
                        var_id = f"{var_id}_combined"
                    variables_found[var_id] = var_info
                
                if combined_vars:
                    st.info(f"🔗 {len(combined_vars)} variables adicionales con patrón combinado")
            
            st.session_state.variables = variables_found
            
            if variables_found:
                st.success(f"✅ Se detectaron {len(variables_found)} variables únicas")
                list_vars = [v for v in variables_found.values() if v['tipo'] == 'lista']
                if list_vars:
                    st.info(f"ℹ️ {len(list_vars)} variables tipo lista detectadas")
            else:
                st.warning("⚠️ No se detectaron variables")
        
        # Paso 3: Configurar variables
        if st.session_state.variables:
            st.markdown("---")
            st.header("3️⃣ Configurar Variables")
            
            # Fusionar variables
            with st.expander("🔀 Fusionar Variables", expanded=False):
                col1, col2 = st.columns([3, 1])
                
                with col1:
                    for var_id in st.session_state.variables.keys():
                        is_selected = var_id in st.session_state.selected_for_merge
                        if st.checkbox(f"`{var_id}`", value=is_selected, key=f"merge_{var_id}"):
                            if var_id not in st.session_state.selected_for_merge:
                                st.session_state.selected_for_merge.append(var_id)
                        else:
                            if var_id in st.session_state.selected_for_merge:
                                st.session_state.selected_for_merge.remove(var_id)
                
                with col2:
                    if len(st.session_state.selected_for_merge) >= 2:
                        new_name = st.text_input("Nombre:", value="_".join(st.session_state.selected_for_merge[:2]), key="merge_name")
                        if st.button("✅ Fusionar", key="merge_btn"):
                            merge_variables(st.session_state.selected_for_merge, new_name)
                    else:
                        st.info("Mínimo 2 variables")
            
            # Resumen
            enabled = sum(1 for v in st.session_state.variables.values() if not v.get('disabled', False))
            disabled = len(st.session_state.variables) - enabled
            st.write(f"**Total: {len(st.session_state.variables)} | ✅ Activas: {enabled} | 🗑️ Desactivadas: {disabled}**")
            
            # Editor de variables
            for idx, (var_id, var_info) in enumerate(st.session_state.variables.items()):
                is_disabled = var_info.get('disabled', False)
                
                if is_disabled:
                    label = f"~~Variable {idx+1}: `{var_id}`~~ 🗑️ DESACTIVADA"
                else:
                    label = f"Variable {idx+1}: `{var_id}`"
                
                with st.expander(label, expanded=False):
                    # Botón desactivar/activar
                    if is_disabled:
                        if st.button("✅ Reactivar", key=f"en_{var_id}"):
                            toggle_variable_enabled(var_id)
                        st.warning("⚠️ Variable desactivada (no aparecerá en YAML)")
                        continue
                    else:
                        if st.button("🗑️ Desactivar", key=f"dis_{var_id}"):
                            toggle_variable_enabled(var_id)
                    
                    st.markdown("---")
                    
                    # 🆕 CONTEXTO (CON CHECKBOX EN LUGAR DE EXPANDER)
                    show_context = st.checkbox("📍 Ver contexto", key=f"ctx_{var_id}")
                    if show_context:
                        try:
                            if file_extension == 'docx':
                                contexts = PatternDetector.extract_variable_context(doc, var_info['original_text'], 20)
                            else:
                                contexts = PatternDetector.extract_variable_context_pptx(prs, var_info['original_text'], 20)

                            if contexts:
                                for i, ctx in enumerate(contexts):
                                    st.markdown(f"**{i+1}.** ({ctx['location']}): `{ctx['before']}`**`{ctx['variable']}`**`{ctx['after']}`")

                                # 🆕 DIVIDIR POR CONTEXTO
                                if len(contexts) > 1:
                                    st.markdown("---")
                                    st.markdown("**✂️ Dividir por contexto:**")

                                    # Inicializar estado para grupos de contexto
                                    if f'context_groups_{var_id}' not in st.session_state:
                                        st.session_state[f'context_groups_{var_id}'] = []

                                    # Botón para añadir grupo
                                    if st.button("➕ Añadir grupo", key=f"add_group_{var_id}"):
                                        st.session_state[f'context_groups_{var_id}'].append({
                                            'name': '',
                                            'contexts': []
                                        })
                                        st.rerun()

                                    # Mostrar grupos existentes
                                    groups = st.session_state[f'context_groups_{var_id}']
                                    if groups:
                                        for group_idx, group in enumerate(groups):
                                            with st.container():
                                                st.markdown(f"**Grupo {group_idx + 1}:**")

                                                gcol1, gcol2 = st.columns([2, 1])
                                                with gcol1:
                                                    group['name'] = st.text_input(
                                                        "Nombre:",
                                                        value=group['name'],
                                                        key=f"gname_{var_id}_{group_idx}",
                                                        placeholder="ej: mes_de_cierre"
                                                    )

                                                with gcol2:
                                                    if st.button("🗑️", key=f"del_group_{var_id}_{group_idx}"):
                                                        st.session_state[f'context_groups_{var_id}'].pop(group_idx)
                                                        st.rerun()

                                                # Seleccionar contextos para este grupo
                                                selected_contexts = []
                                                for ctx_idx in range(len(contexts)):
                                                    is_selected = ctx_idx in group['contexts']
                                                    if st.checkbox(
                                                        f"Contexto {ctx_idx + 1}: {contexts[ctx_idx]['location']}",
                                                        value=is_selected,
                                                        key=f"ctx_sel_{var_id}_{group_idx}_{ctx_idx}"
                                                    ):
                                                        selected_contexts.append(ctx_idx)

                                                group['contexts'] = selected_contexts
                                                st.markdown("---")

                                        # Botón para ejecutar división
                                        if len(groups) >= 2:
                                            if st.button("✨ Dividir variable", key=f"exec_split_{var_id}", type="primary"):
                                                # Validar que todos los grupos tengan nombre y contextos
                                                valid = True
                                                contexts_dict = {}

                                                for group in groups:
                                                    if not group['name'].strip():
                                                        st.error("Todos los grupos deben tener nombre")
                                                        valid = False
                                                        break
                                                    if not group['contexts']:
                                                        st.error("Todos los grupos deben tener al menos un contexto")
                                                        valid = False
                                                        break
                                                    contexts_dict[group['name']] = group['contexts']

                                                if valid:
                                                    split_variable_by_context(var_id, contexts_dict)
                                        else:
                                            st.info("Añade al menos 2 grupos para dividir")
                            else:
                                st.info("No se encontraron apariciones")
                        except Exception as e:
                            st.warning(f"No se pudo extraer contexto: {e}")
                    
                    st.markdown("---")
                    
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        st.text_input("Original:", value=var_info['original_text'], disabled=True, key=f"orig_{var_id}")
                        st.text_input("Nombre:", value=var_id, key=f"name_{var_id}")
                        
                        # División por delimitador
                        st.markdown("**✂️ Dividir por delimitador:**")
                        ca, cb = st.columns(2)
                        with ca:
                            delim = st.text_input("Delim:", "/", key=f"delim_{var_id}", max_chars=3)
                        with cb:
                            st.write("")
                            if st.button("✂️", key=f"split_{var_id}"):
                                split_variable(var_id, delim)
                        
                        # 🆕 División libre
                        st.markdown("**✂️ División libre:**")
                        st.caption(f"Texto: `{var_info['original_text']}` ({len(var_info['original_text'])} chars)")
                        fc1, fc2 = st.columns(2)
                        with fc1:
                            start = st.number_input("Desde:", 0, len(var_info['original_text']), 0, key=f"fs_{var_id}")
                        with fc2:
                            end = st.number_input("Hasta:", 0, len(var_info['original_text']), min(5, len(var_info['original_text'])), key=f"fe_{var_id}")
                        
                        if start < end:
                            prev = var_info['original_text'][start:end]
                            st.info(f"📌 `{prev}`")
                            fname = st.text_input("Nombre:", VariableNormalizer.normalize_name(prev), key=f"fn_{var_id}")
                            if st.button("✨ Crear", key=f"fsplit_{var_id}"):
                                split_variable_free(var_id, start, end, fname)
                    
                    with col2:
                        TIPOS = ['texto', 'numero', 'fecha', 'hora', 'email', 'telefono', 'lista', 'moneda']
                        tipo = st.selectbox("Tipo:", TIPOS, index=TIPOS.index(var_info['tipo']) if var_info['tipo'] in TIPOS else 0, key=f"tipo_{var_id}")
                        st.session_state.variables[var_id]['tipo'] = tipo
                        
                        if tipo == 'lista':
                            opts = st.text_area("Opciones:", "\n".join(var_info.get('opciones', [])), key=f"opts_{var_id}", height=100)
                            st.session_state.variables[var_id]['opciones'] = [o.strip() for o in opts.split('\n') if o.strip()]
                            if var_info.get('opciones'):
                                st.success(f"✨ {len(var_info['opciones'])} opciones")
                        
                        elif tipo == 'moneda':
                            mon = st.selectbox("Moneda:", ['EUR', 'USD'], format_func=lambda x: f"{x} ({'€' if x=='EUR' else '$'})", key=f"mon_{var_id}")
                            meta = st.session_state.variables[var_id].get('meta', {})
                            meta['currency'] = mon
                            st.session_state.variables[var_id]['meta'] = meta
                        
                        elif tipo == 'telefono':
                            PREF = [('España', '+34'), ('Francia', '+33'), ('Portugal', '+351'), ('Italia', '+39'), ('Alemania', '+49'), 
                                    ('Reino Unido', '+44'), ('USA', '+1'), ('México', '+52'), ('Argentina', '+54')]
                            def_idx = 0
                            prev_meta = st.session_state.variables[var_id].get('meta', {})
                            if 'country_code' in prev_meta:
                                for i, (_, c) in enumerate(PREF):
                                    if c == prev_meta['country_code']:
                                        def_idx = i
                                        break
                            sel = st.selectbox("Prefijo:", range(len(PREF)), format_func=lambda i: f"{PREF[i][0]} ({PREF[i][1]})", index=def_idx, key=f"tel_{var_id}")
                            meta = st.session_state.variables[var_id].get('meta', {})
                            meta['country_code'] = PREF[sel][1]
                            st.session_state.variables[var_id]['meta'] = meta
                    
                    with col3:
                        preg = st.text_input("Pregunta:", var_info.get('pregunta', ''), key=f"preg_{var_id}")
                        st.session_state.variables[var_id]['pregunta'] = preg
                        auto = VariableNormalizer.generate_default_question(var_id, tipo)
                        st.info(f"**Auto:** {auto}")
            
            # Generar plantilla
            st.markdown("---")
            st.header("4️⃣ Generar Plantilla")
            
            col1, col2 = st.columns(2)
            with col1:
                nombre = st.text_input("Nombre:", f"plantilla_{uploaded_file.name.rsplit('.', 1)[0]}")
            with col2:
                st.write("")
                st.write("")
                gen_btn = st.button("🚀 Generar", type="primary")
            
            if gen_btn:
                active = {k: v for k, v in st.session_state.variables.items() if not v.get('disabled', False)}
                
                if not active:
                    st.error("❌ No hay variables activas")
                else:
                    with st.spinner("Generando..."):
                        try:
                            work_dir = Path(st.session_state.work_dir)
                            work_dir.mkdir(exist_ok=True)
                            
                            if file_extension == 'docx':
                                template = create_template_docx(st.session_state.original_doc, st.session_state.variables)
                                path = work_dir / f"{nombre}.docx"
                                template.save(str(path))
                            else:
                                template = create_template_pptx(st.session_state.original_doc, st.session_state.variables)
                                path = work_dir / f"{nombre}.pptx"
                                template.save(str(path))
                            
                            yaml_cfg = YAMLManager.create_variable_config(st.session_state.variables)
                            yaml_path = work_dir / f"{nombre}_config.yaml"
                            YAMLManager.save_yaml(yaml_cfg, str(yaml_path))
                            
                            st.success("✅ Generado correctamente")
                            
                            c1, c2 = st.columns(2)
                            with c1:
                                with open(path, 'rb') as f:
                                    st.download_button(f"📄 Plantilla (.{file_extension})", f.read(), f"{nombre}.{file_extension}")
                            with c2:
                                with open(yaml_path, 'rb') as f:
                                    st.download_button("⚙️ Config (YAML)", f.read(), f"{nombre}_config.yaml")
                            
                            st.info(f"**Variables:** {len(st.session_state.variables)} | **Activas:** {len(active)}")
                        except Exception as e:
                            st.error(f"❌ Error: {e}")
                            st.exception(e)


def extract_text_from_docx(doc: Document) -> str:
    """Extrae texto de Word"""
    parts = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for t in doc.tables:
        for r in t.rows:
            for c in r.cells:
                parts.append(c.text)
    for s in doc.sections:
        if s.header:
            for p in s.header.paragraphs:
                parts.append(p.text)
        if s.footer:
            for p in s.footer.paragraphs:
                parts.append(p.text)
    return " ".join(parts)


def extract_text_from_pptx(prs: Presentation) -> str:
    """Extrae texto de PowerPoint"""
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return " ".join(parts)


def create_template_docx(original_doc: Document, variables: dict) -> Document:
    """Crea plantilla Word"""
    doc = original_doc
    active = {k: v for k, v in variables.items() if not v.get('disabled', False)}

    # Inicializar contador de apariciones para variables con contextos específicos
    occurrence_counters = {vid: 0 for vid, vinfo in active.items() if 'context_indices' in vinfo}

    for p in doc.paragraphs:
        replace_in_paragraph(p, active, occurrence_counters)
    for t in doc.tables:
        for r in t.rows:
            for c in r.cells:
                for p in c.paragraphs:
                    replace_in_paragraph(p, active, occurrence_counters)
    for s in doc.sections:
        if s.header:
            for p in s.header.paragraphs:
                replace_in_paragraph(p, active, occurrence_counters)
        if s.footer:
            for p in s.footer.paragraphs:
                replace_in_paragraph(p, active, occurrence_counters)
    return doc


def replace_in_paragraph(paragraph, variables, occurrence_counters=None):
    full = paragraph.text
    new = full

    if occurrence_counters is None:
        occurrence_counters = {}

    for vid, vinfo in variables.items():
        if vinfo.get('disabled', False):
            continue

        orig = vinfo['original_text'] or ""
        clean = clean_pattern_markers(orig)
        placeholder = f"{{{{{vid}}}}}"
        pname = (vinfo.get('pattern') or "").lower()

        # Verificar si esta variable tiene contextos específicos
        context_indices = vinfo.get('context_indices')

        # Buscar texto a reemplazar según el patrón
        search_text = None
        case_insensitive = False

        if pname in ('llaves_simples', 'llaves_dobles', 'corchetes_simples', 'corchetes_dobles'):
            pattern_map = {
                'llaves_simples': f"{{{clean}}}",
                'llaves_dobles': f"{{{{{clean}}}}}",
                'corchetes_simples': f"[{clean}]",
                'corchetes_dobles': f"[[{clean}]]",
            }
            search_text = pattern_map.get(pname)
        elif pname == 'date_with_de':
            # Para patrones de fecha con "de", buscar case-insensitive
            search_text = orig
            case_insensitive = True
        elif orig:
            search_text = orig

        if not search_text:
            continue

        # Verificar si el texto existe (case-insensitive si aplica)
        if case_insensitive:
            if search_text.lower() not in new.lower():
                continue
        else:
            if search_text not in new:
                continue

        # Si tiene contextos específicos, reemplazar selectivamente
        if context_indices is not None:
            # Encontrar todas las apariciones en este párrafo
            positions = []
            start_pos = 0

            if case_insensitive:
                # Búsqueda case-insensitive
                new_lower = new.lower()
                search_lower = search_text.lower()
                while True:
                    pos = new_lower.find(search_lower, start_pos)
                    if pos == -1:
                        break
                    positions.append(pos)
                    start_pos = pos + 1
            else:
                while True:
                    pos = new.find(search_text, start_pos)
                    if pos == -1:
                        break
                    positions.append(pos)
                    start_pos = pos + 1

            # Reemplazar de atrás hacia adelante para mantener índices
            for pos in reversed(positions):
                current_occurrence = occurrence_counters.get(vid, 0)
                if current_occurrence in context_indices:
                    # Determinar longitud del texto a reemplazar
                    if case_insensitive:
                        # Encontrar la longitud exacta del texto en el documento
                        actual_length = len(search_text)
                    else:
                        actual_length = len(search_text)
                    new = new[:pos] + placeholder + new[pos + actual_length:]
                occurrence_counters[vid] = current_occurrence + 1
        else:
            # Reemplazar todas las apariciones
            if pname in ('llaves_simples', 'llaves_dobles', 'corchetes_simples', 'corchetes_dobles'):
                rmap = {
                    'llaves_simples': r'\{' + re.escape(clean) + r'\}',
                    'llaves_dobles': r'\{\{' + re.escape(clean) + r'\}\}',
                    'corchetes_simples': r'\[' + re.escape(clean) + r'\]',
                    'corchetes_dobles': r'\[\[' + re.escape(clean) + r'\]\]',
                }
                new = re.sub(rmap[pname], placeholder, new)
            elif case_insensitive:
                # Reemplazo case-insensitive
                new = re.sub(re.escape(search_text), placeholder, new, flags=re.IGNORECASE)
            else:
                new = new.replace(search_text, placeholder)

    new = sanitize_placeholders(new)

    if new != full:
        for run in paragraph.runs:
            run.text = ''
        if paragraph.runs:
            paragraph.runs[0].text = new
        else:
            paragraph.add_run(new)


def create_template_pptx(original_prs: Presentation, variables: dict) -> Presentation:
    """Crea plantilla PowerPoint"""
    prs = original_prs
    active = {k: v for k, v in variables.items() if not v.get('disabled', False)}

    # Inicializar contador de apariciones para variables con contextos específicos
    occurrence_counters = {vid: 0 for vid, vinfo in active.items() if 'context_indices' in vinfo}

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                replace_in_textframe(shape.text_frame, active, occurrence_counters)
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_in_textframe(cell.text_frame, active, occurrence_counters)
    return prs


def replace_in_textframe(tf, variables, occurrence_counters=None):
    if occurrence_counters is None:
        occurrence_counters = {}

    for para in tf.paragraphs:
        full = ''.join(r.text for r in para.runs)
        new = full

        for vid, vinfo in variables.items():
            if vinfo.get('disabled', False):
                continue

            orig = vinfo['original_text'] or ""
            clean = clean_pattern_markers(orig)
            placeholder = f"{{{{{vid}}}}}"
            pname = (vinfo.get('pattern') or "").lower()

            # Verificar si esta variable tiene contextos específicos
            context_indices = vinfo.get('context_indices')

            # Buscar texto a reemplazar según el patrón
            search_text = None
            case_insensitive = False

            if pname in ('llaves_simples', 'llaves_dobles', 'corchetes_simples', 'corchetes_dobles'):
                pattern_map = {
                    'llaves_simples': f"{{{clean}}}",
                    'llaves_dobles': f"{{{{{clean}}}}}",
                    'corchetes_simples': f"[{clean}]",
                    'corchetes_dobles': f"[[{clean}]]",
                }
                search_text = pattern_map.get(pname)
            elif pname == 'date_with_de':
                # Para patrones de fecha con "de", buscar case-insensitive
                search_text = orig
                case_insensitive = True
            elif orig:
                search_text = orig

            if not search_text:
                continue

            # Verificar si el texto existe (case-insensitive si aplica)
            if case_insensitive:
                if search_text.lower() not in new.lower():
                    continue
            else:
                if search_text not in new:
                    continue

            # Si tiene contextos específicos, reemplazar selectivamente
            if context_indices is not None:
                # Encontrar todas las apariciones en este text frame
                positions = []
                start_pos = 0

                if case_insensitive:
                    # Búsqueda case-insensitive
                    new_lower = new.lower()
                    search_lower = search_text.lower()
                    while True:
                        pos = new_lower.find(search_lower, start_pos)
                        if pos == -1:
                            break
                        positions.append(pos)
                        start_pos = pos + 1
                else:
                    while True:
                        pos = new.find(search_text, start_pos)
                        if pos == -1:
                            break
                        positions.append(pos)
                        start_pos = pos + 1

                # Reemplazar de atrás hacia adelante para mantener índices
                for pos in reversed(positions):
                    current_occurrence = occurrence_counters.get(vid, 0)
                    if current_occurrence in context_indices:
                        actual_length = len(search_text)
                        new = new[:pos] + placeholder + new[pos + actual_length:]
                    occurrence_counters[vid] = current_occurrence + 1
            else:
                # Reemplazar todas las apariciones
                if pname in ('llaves_simples', 'llaves_dobles', 'corchetes_simples', 'corchetes_dobles'):
                    rmap = {
                        'llaves_simples': r'\{' + re.escape(clean) + r'\}',
                        'llaves_dobles': r'\{\{' + re.escape(clean) + r'\}\}',
                        'corchetes_simples': r'\[' + re.escape(clean) + r'\]',
                        'corchetes_dobles': r'\[\[' + re.escape(clean) + r'\]\]',
                    }
                    new = re.sub(rmap[pname], placeholder, new)
                elif case_insensitive:
                    # Reemplazo case-insensitive
                    new = re.sub(re.escape(search_text), placeholder, new, flags=re.IGNORECASE)
                else:
                    new = new.replace(search_text, placeholder)

        new = sanitize_placeholders(new)

        if new != full:
            for run in para.runs:
                run.text = ''
            if para.runs:
                para.runs[0].text = new
            else:
                para.add_run(new)


if __name__ == "__main__":
    main()